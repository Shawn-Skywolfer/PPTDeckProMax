"""
PPT Deck Pro Max — Optional Backend
FastAPI server for PPTX generation, screenshot capture, and montage build.
Auto-detects optional dependencies (python-pptx, playwright, PIL).
"""

import os
import sys
import json
import io
import tempfile
from typing import Optional, List, Dict, Any
from datetime import datetime

# FastAPI
try:
    from fastapi import FastAPI, HTTPException, UploadFile, File, Form
    from fastapi.responses import StreamingResponse, JSONResponse
    from fastapi.middleware.cors import CORSMiddleware
except ImportError:
    print("ERROR: fastapi not installed. Run: pip install fastapi uvicorn")
    sys.exit(1)

# python-pptx (optional)
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("WARNING: python-pptx not installed. PPTX generation disabled.")

# Pillow (optional)
try:
    from PIL import Image, ImageDraw, ImageFont
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    print("WARNING: Pillow not installed. Montage generation disabled.")

# Playwright (optional)
try:
    from playwright.async_api import async_playwright
    PLAYWRIGHT_AVAILABLE = True
except ImportError:
    PLAYWRIGHT_AVAILABLE = False
    print("WARNING: playwright not installed. Screenshot capture disabled.")

app = FastAPI(title="PPT Deck Pro Max Backend", version="1.0.0")

# CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:*", "http://127.0.0.1:*", "null"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.get("/api/health")
async def health_check():
    return {
        "status": "ok",
        "capabilities": {
            "pptx": PPTX_AVAILABLE,
            "screenshot": PLAYWRIGHT_AVAILABLE,
            "montage": PIL_AVAILABLE
        }
    }


@app.post("/api/build-pptx")
async def build_pptx(request: Dict[str, Any]):
    """Build PPTX from deck structure JSON."""
    if not PPTX_AVAILABLE:
        raise HTTPException(status_code=501, detail="python-pptx not installed")

    try:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        theme = request.get("theme_tokens", {})
        colors = theme.get("colors", {})
        fonts = theme.get("fonts", {})

        # Default colors
        primary = _parse_color(colors.get("primary", "#6366f1"))
        secondary = _parse_color(colors.get("secondary", "#8b5cf6"))
        text_dark = _parse_color(colors.get("text_dark", "#1e293b"))
        text_light = _parse_color(colors.get("text_light", "#64748b"))
        bg = _parse_color(colors.get("bg", "#ffffff"))
        accent = _parse_color(colors.get("accent", "#f59e0b"))

        slides = request.get("slides", [])

        for slide_data in slides:
            layout_type = slide_data.get("layout", "content")
            title = slide_data.get("title", "")
            content = slide_data.get("content_html", "")

            slide_layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(slide_layout)

            if layout_type == "title":
                _add_title_slide(slide, title, content, primary, text_dark)
            elif layout_type == "content":
                _add_content_slide(slide, title, content, text_dark, text_light)
            elif layout_type == "two_column":
                _add_two_column_slide(slide, title, content, text_dark, text_light)
            elif layout_type == "chart":
                _add_chart_slide(slide, title, slide_data.get("chart_data", {}), text_dark, primary)
            else:
                _add_content_slide(slide, title, content, text_dark, text_light)

        # Save to memory
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        filename = request.get("project_id", "deck") + ".pptx"
        return StreamingResponse(
            output,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/capture-screenshot")
async def capture_screenshot(request: Dict[str, Any]):
    """Capture screenshot from HTML."""
    if not PLAYWRIGHT_AVAILABLE:
        raise HTTPException(status_code=501, detail="playwright not installed")

    try:
        html = request.get("html", "")
        width = request.get("width", 1280)
        height = request.get("height", 960)

        async with async_playwright() as p:
            browser = await p.chromium.launch()
            page = await browser.new_page(viewport={"width": width, "height": height})
            await page.set_content(html)
            await page.wait_for_timeout(1000)
            screenshot = await page.screenshot(type="png")
            await browser.close()

        return StreamingResponse(
            io.BytesIO(screenshot),
            media_type="image/png"
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/build-montage")
async def build_montage(
    images: List[UploadFile] = File(...),
    layout: str = Form("{\"rows\": 2, \"cols\": 3, \"padding\": 10}")
):
    """Build montage grid from images."""
    if not PIL_AVAILABLE:
        raise HTTPException(status_code=501, detail="Pillow not installed")

    try:
        layout_cfg = json.loads(layout)
        rows = layout_cfg.get("rows", 2)
        cols = layout_cfg.get("cols", 3)
        padding = layout_cfg.get("padding", 10)

        thumb_w = 320
        thumb_h = 240
        montage_w = cols * thumb_w + (cols + 1) * padding
        montage_h = rows * thumb_h + (rows + 1) * padding

        montage = Image.new("RGB", (montage_w, montage_h), (245, 245, 245))

        for idx, img_file in enumerate(images):
            if idx >= rows * cols:
                break
            content = await img_file.read()
            img = Image.open(io.BytesIO(content))
            img = img.resize((thumb_w, thumb_h), Image.LANCZOS)

            row = idx // cols
            col = idx % cols
            x = col * thumb_w + (col + 1) * padding
            y = row * thumb_h + (row + 1) * padding
            montage.paste(img, (x, y))

        output = io.BytesIO()
        montage.save(output, format="PNG")
        output.seek(0)

        return StreamingResponse(output, media_type="image/png")

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ===== Helpers =====

def _parse_color(hex_color: str) -> RGBColor:
    """Parse hex color to RGBColor."""
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 3:
        hex_color = "".join(c * 2 for c in hex_color)
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def _add_title_slide(slide, title, subtitle, primary, text_dark):
    """Add a title slide."""
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide.shapes._spTree.getparent().getparent().sldSz.cx, slide.shapes._spTree.getparent().getparent().sldSz.cy)
    bg.fill.solid()
    bg.fill.fore_color.rgb = primary
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8.4), Inches(1.5))
    tf = title_box.text_frame
    tf.text = title
    _set_text_style(tf, font_size=44, color=RGBColor(255, 255, 255), bold=True)

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(8.4), Inches(0.8))
    tf = sub_box.text_frame
    tf.text = subtitle
    _set_text_style(tf, font_size=24, color=RGBColor(200, 200, 200))


def _add_content_slide(slide, title, content, text_dark, text_light):
    """Add a content slide."""
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = title
    _set_text_style(tf, font_size=28, color=text_dark, bold=True)

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    tf = content_box.text_frame
    tf.text = content
    _set_text_style(tf, font_size=14, color=text_dark)
    tf.word_wrap = True


def _add_two_column_slide(slide, title, content, text_dark, text_light):
    """Add a two-column slide."""
    _add_content_slide(slide, title, content, text_dark, text_light)


def _add_chart_slide(slide, title, chart_data, text_dark, primary):
    """Add a chart slide (placeholder)."""
    _add_content_slide(slide, title, "[Chart Placeholder]", text_dark, text_dark)


def _set_text_style(text_frame, font_name="Microsoft YaHei", font_size=14, color=RGBColor(0, 0, 0), bold=False, align=PP_ALIGN.LEFT):
    """Apply consistent text styling."""
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = bold


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8765)
