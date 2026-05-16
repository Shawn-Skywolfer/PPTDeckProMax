#!/usr/bin/env python3
"""
Build PPT Deck: 全球能源格局重构 — 霍尔木兹危机后的新能源机遇
Based on report_detailed.md using ppt-deck-pro-max skill workflow
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
import os

# ============================================================
# COLOR SYSTEM
# ============================================================
DEEP_BLUE = RGBColor(10, 37, 64)      # #0A2540
CRISIS_RED = RGBColor(230, 57, 70)    # #E63946
ENERGY_GOLD = RGBColor(255, 183, 3)   # #FFB703
GRAY_BLUE = RGBColor(74, 108, 140)    # #4A6C8C
LIGHT_GRAY = RGBColor(245, 245, 245)  # #F5F5F5
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(22, 163, 74)         # #16A34A
AMBER = RGBColor(202, 138, 4)         # #CA8A04
ORANGE = RGBColor(234, 88, 12)        # #EA580C
DARK_RED = RGBColor(220, 38, 38)      # #DC2626
BLACK = RGBColor(0, 0, 0)

def set_text_frame_style(text_frame, font_name="Microsoft YaHei", font_size=14, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    """Apply consistent text styling"""
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = bold


def _style_chart_series(series, color):
    """Apply color and data labels to a chart series."""
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color
    series.format.line.fill.background()
    dl = series.data_labels
    dl.show_value = True
    dl.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    dl.font.size = Pt(10)
    dl.font.name = "Microsoft YaHei"
    dl.font.color.rgb = BLACK


def _style_chart_axes(chart, text_color=GRAY_BLUE):
    """Style chart category and value axes."""
    for axis in (chart.category_axis, chart.value_axis):
        axis.tick_labels.font.size = Pt(10)
        axis.tick_labels.font.name = "Microsoft YaHei"
        axis.tick_labels.font.color.rgb = text_color
        axis.has_major_gridlines = False
    chart.value_axis.has_minor_gridlines = False


def add_title_slide(prs, title, subtitle, date_str):
    """P1: Cover slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background shape
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DEEP_BLUE
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8.4), Inches(1.5))
    tf = title_box.text_frame
    tf.text = title
    set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=44, color=ENERGY_GOLD, bold=True, align=PP_ALIGN.LEFT)

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(8.4), Inches(0.8))
    tf = sub_box.text_frame
    tf.text = subtitle
    set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=24, color=WHITE, bold=False, align=PP_ALIGN.LEFT)

    # Date
    date_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(8.4), Inches(0.5))
    tf = date_box.text_frame
    tf.text = date_str
    set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=16, color=GRAY_BLUE, bold=False, align=PP_ALIGN.LEFT)

    # Decorative line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.0), Inches(2.0), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = CRISIS_RED
    line.line.fill.background()

    return slide

def add_two_section_slide(prs, title, conclusion, left_content, right_content, footer=""):
    """Generic two-section slide layout"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = title
    set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=28, color=DEEP_BLUE, bold=True, align=PP_ALIGN.LEFT)

    # Conclusion bar
    conc_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(9), Inches(0.4))
    tf = conc_box.text_frame
    tf.text = conclusion
    set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=14, color=CRISIS_RED, bold=True, align=PP_ALIGN.LEFT)

    # Left section
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.3), Inches(5.5))
    tf = left_box.text_frame
    tf.text = left_content
    set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=12, color=BLACK, bold=False, align=PP_ALIGN.LEFT)
    tf.word_wrap = True

    # Right section
    right_box = slide.shapes.add_textbox(Inches(5.0), Inches(1.5), Inches(4.3), Inches(5.5))
    tf = right_box.text_frame
    tf.text = right_content
    set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=12, color=BLACK, bold=False, align=PP_ALIGN.LEFT)
    tf.word_wrap = True

    # Footer
    if footer:
        foot_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
        tf = foot_box.text_frame
        tf.text = footer
        set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=9, color=GRAY_BLUE, bold=False, align=PP_ALIGN.LEFT)

    return slide

def add_bar_chart_slide(prs, title, conclusion, categories, values, series_name="数据", footer=""):
    """Add a slide with bar chart"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = title
    set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=28, color=DEEP_BLUE, bold=True, align=PP_ALIGN.LEFT)

    # Conclusion
    conc_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(9), Inches(0.4))
    tf = conc_box.text_frame
    tf.text = conclusion
    set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=14, color=CRISIS_RED, bold=True, align=PP_ALIGN.LEFT)

    # Chart data
    chart_data = ChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)

    # Add chart
    x, y, cx, cy = Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.0)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(10)
    chart.legend.font.name = "Microsoft YaHei"

    # Style series and axes
    _style_chart_series(chart.series[0], CRISIS_RED)
    _style_chart_axes(chart)

    # Footer
    if footer:
        foot_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
        tf = foot_box.text_frame
        tf.text = footer
        set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=9, color=GRAY_BLUE, bold=False, align=PP_ALIGN.LEFT)

    return slide

def add_big_number_slide(prs, title, conclusion, numbers, labels, footer=""):
    """Slide with big numbers and cards"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = title
    set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=28, color=DEEP_BLUE, bold=True, align=PP_ALIGN.LEFT)

    # Conclusion
    conc_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(9), Inches(0.4))
    tf = conc_box.text_frame
    tf.text = conclusion
    set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=14, color=CRISIS_RED, bold=True, align=PP_ALIGN.LEFT)

    # Number cards
    n = len(numbers)
    card_width = 8.5 / n
    for i, (num, label) in enumerate(zip(numbers, labels)):
        x = Inches(0.5 + i * card_width)
        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0), Inches(card_width - 0.2), Inches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = GRAY_BLUE

        # Number
        num_box = slide.shapes.add_textbox(x, Inches(2.3), Inches(card_width - 0.2), Inches(1.0))
        tf = num_box.text_frame
        tf.text = str(num)
        set_text_frame_style(tf, font_name="DIN Alternate", font_size=36, color=CRISIS_RED, bold=True, align=PP_ALIGN.CENTER)

        # Label
        label_box = slide.shapes.add_textbox(x, Inches(3.3), Inches(card_width - 0.2), Inches(0.8))
        tf = label_box.text_frame
        tf.text = label
        set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=12, color=GRAY_BLUE, bold=False, align=PP_ALIGN.CENTER)
        tf.word_wrap = True

    # Footer
    if footer:
        foot_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
        tf = foot_box.text_frame
        tf.text = footer
        set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=9, color=GRAY_BLUE, bold=False, align=PP_ALIGN.LEFT)

    return slide

def add_comparison_slide(prs, title, conclusion, left_title, left_data, right_title, right_data, footer=""):
    """Side-by-side comparison slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = title
    set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=28, color=DEEP_BLUE, bold=True, align=PP_ALIGN.LEFT)

    # Conclusion
    conc_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(9), Inches(0.4))
    tf = conc_box.text_frame
    tf.text = conclusion
    set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=14, color=CRISIS_RED, bold=True, align=PP_ALIGN.LEFT)

    # Left side background
    left_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.5), Inches(4.5), Inches(5.0))
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = RGBColor(240, 253, 244)  # Light green
    left_bg.line.color.rgb = GREEN

    # Left title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(4.1), Inches(0.5))
    tf = left_title_box.text_frame
    tf.text = left_title
    set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=20, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

    # Left data
    left_data_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(4.1), Inches(3.8))
    tf = left_data_box.text_frame
    tf.text = left_data
    set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=13, color=BLACK, bold=False, align=PP_ALIGN.LEFT)
    tf.word_wrap = True

    # VS label
    vs_box = slide.shapes.add_textbox(Inches(4.6), Inches(3.5), Inches(0.8), Inches(0.6))
    tf = vs_box.text_frame
    tf.text = "VS"
    set_text_frame_style(tf, font_name="DIN Alternate", font_size=32, color=GRAY_BLUE, bold=True, align=PP_ALIGN.CENTER)

    # Right side background
    right_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.5), Inches(4.5), Inches(5.0))
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = RGBColor(254, 242, 242)  # Light red
    right_bg.line.color.rgb = DARK_RED

    # Right title
    right_title_box = slide.shapes.add_textbox(Inches(5.4), Inches(1.7), Inches(4.1), Inches(0.5))
    tf = right_title_box.text_frame
    tf.text = right_title
    set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=20, color=DARK_RED, bold=True, align=PP_ALIGN.CENTER)

    # Right data
    right_data_box = slide.shapes.add_textbox(Inches(5.4), Inches(2.3), Inches(4.1), Inches(3.8))
    tf = right_data_box.text_frame
    tf.text = right_data
    set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=13, color=BLACK, bold=False, align=PP_ALIGN.LEFT)
    tf.word_wrap = True

    # Footer
    if footer:
        foot_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
        tf = foot_box.text_frame
        tf.text = footer
        set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=9, color=GRAY_BLUE, bold=False, align=PP_ALIGN.LEFT)

    return slide

def add_timeline_slide(prs, title, conclusion, events, footer=""):
    """Timeline slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = title
    set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=28, color=DEEP_BLUE, bold=True, align=PP_ALIGN.LEFT)

    # Conclusion
    conc_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(9), Inches(0.4))
    tf = conc_box.text_frame
    tf.text = conclusion
    set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=14, color=CRISIS_RED, bold=True, align=PP_ALIGN.LEFT)

    # Timeline line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.5), Inches(8.4), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = GRAY_BLUE
    line.line.fill.background()

    # Events
    n = len(events)
    step = 8.4 / (n - 1) if n > 1 else 4.2
    for i, (date, desc) in enumerate(events):
        x = Inches(0.8 + i * step)
        # Dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x - Inches(0.1), Inches(3.4), Inches(0.2), Inches(0.2))
        dot.fill.solid()
        dot.fill.fore_color.rgb = CRISIS_RED if i < 2 else GRAY_BLUE
        dot.line.fill.background()

        # Date
        date_box = slide.shapes.add_textbox(x - Inches(0.8), Inches(2.8), Inches(1.6), Inches(0.4))
        tf = date_box.text_frame
        tf.text = date
        set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=11, color=CRISIS_RED if i < 2 else GRAY_BLUE, bold=True, align=PP_ALIGN.CENTER)

        # Description
        desc_box = slide.shapes.add_textbox(x - Inches(0.9), Inches(3.8), Inches(1.8), Inches(1.5))
        tf = desc_box.text_frame
        tf.text = desc
        set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=10, color=BLACK, bold=False, align=PP_ALIGN.CENTER)
        tf.word_wrap = True

    # Footer
    if footer:
        foot_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
        tf = foot_box.text_frame
        tf.text = footer
        set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=9, color=GRAY_BLUE, bold=False, align=PP_ALIGN.LEFT)

    return slide

def add_matrix_slide(prs, title, conclusion, quadrants, highlights, footer=""):
    """2x2 matrix slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = title
    set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=28, color=DEEP_BLUE, bold=True, align=PP_ALIGN.LEFT)

    # Conclusion
    conc_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(9), Inches(0.4))
    tf = conc_box.text_frame
    tf.text = conclusion
    set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=14, color=CRISIS_RED, bold=True, align=PP_ALIGN.LEFT)

    # Matrix background
    matrix_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.8), Inches(7.0), Inches(4.5))
    matrix_bg.fill.solid()
    matrix_bg.fill.fore_color.rgb = LIGHT_GRAY
    matrix_bg.line.color.rgb = GRAY_BLUE

    # Axis labels
    y_label = slide.shapes.add_textbox(Inches(0.2), Inches(3.5), Inches(1.0), Inches(0.5))
    tf = y_label.text_frame
    tf.text = "LNG依存度 ↑"
    set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=11, color=GRAY_BLUE, bold=True, align=PP_ALIGN.CENTER)

    x_label = slide.shapes.add_textbox(Inches(4.0), Inches(6.4), Inches(2.0), Inches(0.4))
    tf = x_label.text_frame
    tf.text = "支付能力 →"
    set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=11, color=GRAY_BLUE, bold=True, align=PP_ALIGN.CENTER)

    # Quadrant labels
    quadrants_pos = [
        (Inches(2.0), Inches(2.2), "潜力市场\n(高依存/低支付)"),
        (Inches(6.0), Inches(2.2), "黄金机会\n(高依存/高支付)"),
        (Inches(2.0), Inches(5.0), "排除\n(低依存/低支付)"),
        (Inches(6.0), Inches(5.0), "观望\n(低依存/高支付)"),
    ]

    for x, y, text in quadrants_pos:
        box = slide.shapes.add_textbox(x, y, Inches(2.0), Inches(0.8))
        tf = box.text_frame
        tf.text = text
        set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=10, color=GRAY_BLUE, bold=False, align=PP_ALIGN.CENTER)

    # Highlights
    for x, y, text, color in highlights:
        box = slide.shapes.add_textbox(x, y, Inches(1.2), Inches(0.5))
        tf = box.text_frame
        tf.text = text
        set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=11, color=color, bold=True, align=PP_ALIGN.CENTER)

    # Footer
    if footer:
        foot_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
        tf = foot_box.text_frame
        tf.text = footer
        set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=9, color=GRAY_BLUE, bold=False, align=PP_ALIGN.LEFT)

    return slide

def add_cta_slide(prs, title, actions, timeline, footer=""):
    """CTA slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DEEP_BLUE
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(1.0))
    tf = title_box.text_frame
    tf.text = title
    set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=40, color=ENERGY_GOLD, bold=True, align=PP_ALIGN.LEFT)

    # Actions
    y_pos = 3.0
    for i, action in enumerate(actions):
        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(y_pos), Inches(0.4), Inches(0.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ENERGY_GOLD
        circle.line.fill.background()

        num_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.05), Inches(0.4), Inches(0.3))
        tf = num_box.text_frame
        tf.text = str(i + 1)
        set_text_frame_style(tf, font_name="DIN Alternate", font_size=16, color=DEEP_BLUE, bold=True, align=PP_ALIGN.CENTER)

        # Action text
        act_box = slide.shapes.add_textbox(Inches(1.4), Inches(y_pos), Inches(7.5), Inches(0.5))
        tf = act_box.text_frame
        tf.text = action
        set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=16, color=WHITE, bold=False, align=PP_ALIGN.LEFT)

        y_pos += 0.7

    # Timeline
    timeline_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(8.4), Inches(0.8))
    tf = timeline_box.text_frame
    tf.text = timeline
    set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=14, color=GRAY_BLUE, bold=False, align=PP_ALIGN.LEFT)

    # Footer
    if footer:
        foot_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(8.4), Inches(0.3))
        tf = foot_box.text_frame
        tf.text = footer
        set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=9, color=GRAY_BLUE, bold=False, align=PP_ALIGN.LEFT)

    return slide

# ============================================================
# BUILD THE DECK
# ============================================================
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

output_path = "D:/AI/AI-PPT/PPT-Deck-Pro-Max/全球能源格局重构_霍尔木兹危机后的新能源机遇.pptx"

# P1: Cover
add_title_slide(prs, "全球能源格局重构", "霍尔木兹危机后的新能源机遇", "2026年5月14日 | 光伏与储能的战略窗口期")

# P2: 霍尔木兹封锁现状
add_big_number_slide(prs,
    "霍尔木兹海峡封锁已延长至5月底",
    "EIA最新假设比此前预期更悲观，每日2100万桶原油运输中断",
    ["21M", "81天", "+48.3%"],
    ["每日通过量\n(barrels)", "封锁时长\n(2/23起)", "Brent涨幅\n(较基准)"],
    "数据来源：EIA, Reuters, OilPrice.com")

# P3: 油价高位整理
add_bar_chart_slide(prs,
    "Brent维持$105上方高位，现货溢价反映供应紧张",
    "原油市场未出现回落迹象，OPEC Basket $115.1显示现货市场继续偏紧",
    ["Brent", "WTI", "OPEC Basket", "Indian Basket"],
    [105.70, 101.20, 115.1, 109.1],
    "当前价格 (USD)",
    "数据来源：OilPrice.com, 2026年5月14日")

# P4: LNG供应链断裂
add_two_section_slide(prs,
    "卡塔尔LNG复产无时间表，全球20%供应中断",
    "Ras Laffan设施受损导致全球LNG供应链最严峻中断，南亚三国首当其冲",
    "【关键数据】\n• 全球LNG供应中断: 20%\n• JKM现货价格: $16.99/MMBtu\n• TTF欧洲基准: ~$15.7/MMBtu\n• Henry Hub美国: $2.871/MMBtu\n• 区域价差: $14.12/MMBtu\n\n【受影响国家】\n• 巴基斯坦: 几乎完全断供\n• 孟加拉国: 严重短缺\n• 印度: 长期合同受限\n• 泰国: 电力成本上升",
    "【事件时间线】\n• 2月23日: 伊朗战争爆发\n• 3月初: 霍尔木兹海峡封锁\n• 3月中旬: 卡塔尔Ras Laffan受损\n• 5月12日: EIA确认封锁延至月底\n• 至今: 复产仍无时间表\n\n【IEA警告】\n• 原油库存以创纪录速度下降\n• 紧张将持续至2030年\n• 全球LNG产量同比-8%",
    "数据来源：Reuters, IEA Gas Market Report Q2-2026, TradingEconomics")

# P5: 全球汽油价格排名
add_bar_chart_slide(prs,
    "全球汽油均价$1.46/L，较基准上涨21.7%",
    "香港$4.16/L领跑，欧洲普遍$2.3+，亚洲分化严重",
    ["香港", "以色列", "丹麦", "荷兰", "希腊", "芬兰", "瑞士", "新加坡", "法国", "挪威"],
    [4.16, 2.86, 2.76, 2.75, 2.46, 2.46, 2.44, 2.42, 2.38, 2.35],
    "汽油价格 (USD/L)",
    "数据来源：GlobalPetrolPrices.com (11-May-2026), Serper验证 | ⚠️ 英国-38.7%基准异常")

# P6: 全球柴油价格排名
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = title_box.text_frame
tf.text = "柴油涨幅超汽油，工业需求支撑韧性"
set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=28, color=DEEP_BLUE, bold=True, align=PP_ALIGN.LEFT)
conc_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(9), Inches(0.4))
tf = conc_box.text_frame
tf.text = "柴油+31.7% vs 汽油+21.7%，溢价扩大至$0.12/L，反映工业和运输需求强劲"
set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=14, color=CRISIS_RED, bold=True, align=PP_ALIGN.LEFT)
chart_data = ChartData()
chart_data.categories = ["香港", "新加坡", "瑞士", "荷兰", "丹麦", "芬兰", "以色列", "英国", "法国", "意大利"]
chart_data.add_series("汽油", [4.16, 2.42, 2.44, 2.75, 2.76, 2.46, 2.86, 2.12, 2.38, 2.27])
chart_data.add_series("柴油", [4.70, 3.27, 2.77, 2.71, 2.74, 2.73, 2.71, 2.55, 2.53, 2.37])
x, y, cx, cy = Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.0)
chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data).chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.font.size = Pt(10)
chart.legend.font.name = "Microsoft YaHei"

# Style both series
_style_chart_series(chart.series[0], GRAY_BLUE)
_style_chart_series(chart.series[1], CRISIS_RED)
_style_chart_axes(chart)

foot_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
tf = foot_box.text_frame
tf.text = "数据来源：GlobalPetrolPrices.com, Serper验证 | ⚠️ 日本-25.2%基准异常"
set_text_frame_style(tf, font_name="Microsoft YaHei", font_size=9, color=GRAY_BLUE, bold=False, align=PP_ALIGN.LEFT)

# P7: LNG依存度排名
add_bar_chart_slide(prs,
    "南亚三国LNG依存度超75%，处于极高风险",
    "巴基斯坦99%、孟加拉国85%、斯里兰卡75%——几乎没有替代方案",
    ["巴基斯坦", "孟加拉国", "斯里兰卡", "印度", "泰国", "韩国", "菲律宾", "日本", "中国", "台湾地区"],
    [99, 85, 75, 60, 55, 20, 17, 11, 10, 8],
    "中东LNG依存度 (%)",
    "数据来源：IEA, IEEFA, Oxford Institute for Energy Studies | $1070亿投资面临风险")

# P8: 能源价格传导
add_two_section_slide(prs,
    "天然气断供推高电力成本，巴基斯坦电价年涨35%",
    "LNG依存度与电价涨幅高度相关，南亚国家承受双重打击",
    "【居民气价 (USD/kWh)】\n• 日本: ~0.095 (+5%)\n• 韩国: ~0.085 (+3%)\n• 中国: ~0.050 (0%)\n• 印度: ~0.035 (+15%)\n• 巴基斯坦: ~0.040 (+50%+)\n• 孟加拉国: ~0.030 (+40%+)\n• 泰国: ~0.045 (+10%)\n• 英国: ~0.080 (-5%)\n• 德国: ~0.090 (-8%)\n• 荷兰: ~0.085 (-10%)",
    "【居民电价 (USD/kWh)】\n• 日本: $0.28 (+8%年涨)\n• 韩国: $0.15 (+10%年涨)\n• 中国: $0.10 (+5%年涨)\n• 印度: $0.10 (+12%年涨)\n• 巴基斯坦: $0.08 (+35%年涨) 🔴\n• 孟加拉国: $0.06 (+25%年涨) 🔴\n• 泰国: $0.12 (+15%年涨)\n• 英国: $0.32 (+5%年涨)\n• 德国: $0.35 (+3%年涨)\n• 荷兰: $0.26 (+2%年涨)\n\n【全球LNG现货】\n• JKM: $16.99/MMBtu\n• TTF: ~$15.7/MMBtu\n• Henry Hub: $2.871/MMBtu",
    "数据来源：GlobalPetrolPrices.com Q1 2026, TradingEconomics, IEA")

# P9: 中国 vs 巴基斯坦
add_comparison_slide(prs,
    "能源自主的两种命运",
    "多元化决定危机中的经济安全——中国通过光伏+核电+管道气将中东依赖降至10%",
    "中国：能源多元化成功者",
    "• 中东LNG依赖: 10%\n• 能源自给率: ~85%\n• 电价年涨: +5%\n• 能源结构:\n  - 光伏: 35%\n  - 核电: 15%\n  - 管道气: 20%\n  - 水电: 20%\n  - 其他: 10%\n• 关键优势:\n  - 俄罗斯管道气进口\n  - 国内煤炭/可再生能源\n  - 能源成本可控",
    "巴基斯坦：能源脆弱典型案例",
    "• 中东LNG依赖: 99%\n• 能源自给率: ~5%\n• 电价年涨: +35%\n• 能源结构:\n  - 卡塔尔LNG: 99%\n  - 其他: 1%\n• 关键困境:\n  - 拒绝$17+/MMBtu投标\n  - 工业停产\n  - 天然气严重短缺\n  - 几乎无多元化空间",
    "数据来源：IEA, 各国能源部公开数据")

# P10: 多元化成功者
add_two_section_slide(prs,
    "核电重启、光伏扩张、管道气——三大有效路径",
    "日本、韩国、中国、欧洲已建立能源缓冲，证明自主能力可以建设",
    "【日本】\n• LNG需求降8%\n• $31亿夏季电力补贴\n• 核电重启中\n\n【韩国】\n• 7GW核电新容量至2033\n• LNG需求预计降20%\n• 已重启煤电应急",
    "【中国】\n• 中东依赖仅10%\n• 俄罗斯管道气+国内生产\n• 光伏装机全球第一\n\n【欧洲】\n• 美国LNG占进口2/3\n• 摆脱卡塔尔依赖\n• 可再生能源加速",
    "数据来源：IEA, 各国能源政策公告")

# P11: 多元化失败者
add_two_section_slide(prs,
    "$1070亿LNG投资面临沉没风险",
    "巴基斯坦、孟加拉国、斯里兰卡几乎无多元化空间，亚洲LNG进口降至六年新低",
    "【三国困境】\n• 巴基斯坦:\n  - 拒绝$17+/MMBtu投标\n  - 工业停产\n  - 能源危机\n• 孟加拉国:\n  - 电力短缺\n  - 财政困难\n  - 补贴难以维持\n• 斯里兰卡:\n  - 全面能源危机\n  - 经济严重受损",
    "【$1070亿投资风险分解】\n• 印度: LNG接收站/管道\n• 巴基斯坦: 终端设施\n• 孟加拉国: 浮式储存\n• 合计: $1070亿\n\n【趋势】\n• 亚洲LNG进口降至六年新低\n• 日本燃气电厂减产\n• 韩国重启煤电\n• 全球LNG产量同比-8%",
    "数据来源：IEEFA, Reuters, IEA")

# P12: 光储一体化
add_two_section_slide(prs,
    "光储一体化：最快的能源自主路径",
    "部署周期6-12个月，不依赖现有电网，适配基础设施薄弱市场",
    "【光储一体化架构】\n• 光伏板发电\n• 逆变器转换\n• 储能电池储存\n• 微网/离网输出\n\n【核心优势】\n✓ 部署快: 6-12个月\n✓ 离网可用: 不依赖电网\n✓ 成本可控: 长期固定成本\n✓ 模块化: 按需扩展",
    "【部署周期对比】\n• 光储微网: 6-12个月 ✅\n• LNG接收终端: 3-5年\n• 核电站: 10年+\n• 跨境管道: 5-8年\n\n【适用场景】\n• 电网薄弱地区\n• 离岛/偏远地区\n• 应急备用电源\n• 工业微网",
    "公司技术白皮书")

# P13: 为什么是我们
add_comparison_slide(prs,
    "离网/微网方案，3-6个月可运行",
    "传统并网光伏需要电网改造，我们的微网方案独立运行",
    "传统并网方案",
    "• 电网评估: 2-3个月\n• 变电站改造: 3-6个月\n• 并网审批: 2-4个月\n• 设备安装: 1-2个月\n• 调试并网: 1-2个月\n\n总计: 9-17个月\n\n❌ 依赖现有电网\n❌ 审批流程复杂\n❌ 基础设施要求高",
    "光储微网方案",
    "• 现场评估: 2-4周\n• 方案设计: 2-4周\n• 设备运输: 4-8周\n• 安装调试: 4-8周\n• 投入运行: 即时\n\n总计: 3-6个月 ✅\n\n✅ 独立运行\n✅ 审批简单\n✅ 适配薄弱电网",
    "公司产品资料")

# P14: IEA预警
add_timeline_slide(prs,
    "天然气紧张将持续至2030年",
    "全球LNG产量同比-8%，这不是短期波动，而是结构性转变",
    [
        ("2026 Q2", "霍尔木兹封锁\n海峡中断\n危机爆发"),
        ("2027", "卡塔尔复产\n(假设)\n供应部分恢复"),
        ("2028", "美国新LNG\n产能上线\n增量供应"),
        ("2029", "可再生能源\n装机加速\n替代效应显现"),
        ("2030", "供需再平衡\n但成本结构\n已永久改变"),
    ],
    "数据来源：IEA Gas Market Report Q2-2026")

# P15: LNG依存国新能源机会矩阵
add_matrix_slide(prs,
    "LNG依存国新能源机会矩阵",
    "依存度 × 支付能力 × 政策友好度 = 黄金机会",
    [],
    [
        (Inches(5.5), Inches(3.0), "泰国", GREEN),
        (Inches(6.5), Inches(2.5), "日本", GREEN),
        (Inches(5.0), Inches(3.5), "印度", AMBER),
        (Inches(6.0), Inches(4.0), "韩国", AMBER),
        (Inches(4.0), Inches(4.5), "巴基斯坦", GRAY_BLUE),
        (Inches(3.5), Inches(4.8), "孟加拉国", GRAY_BLUE),
    ],
    "数据来源：综合IEA、World Bank、各国能源部数据 | 气泡大小=市场规模，颜色=政策友好度")

# P16: TOP3首单入口
add_two_section_slide(prs,
    "TOP3首单入口：泰国试点、日本储能、印度规模",
    "分阶段进入——泰国验证→日本大单→印度规模",
    "【泰国：试点首选】\n• LNG依存度: 55%\n• 政策环境: 友好\n• 支付能力: 中等\n• 优势:\n  - 已取消油价上限\n  - 转向新能源意愿强\n  - 电力成本上升驱动\n\n【进入策略】\nQ2评估 → Q3试点",
    "【日本：储能大单】\n• LNG依存度: 11%\n• 政策环境: 高补贴\n• 支付能力: 高\n• 优势:\n  - $31亿夏季补贴\n  - 核电重启仍有缺口\n  - 储能需求大\n\n【印度：规模市场】\n• LNG依存度: 60%\n• 政策环境: 政府推动\n• 支付能力: 中等\n• 优势:\n  - $1070亿投资风险\n  - 政府大力推动光伏\n  - 规模巨大\n\n【进入策略】\nQ4洽谈日本 → 2027 Q1投标印度",
    "各国政策公告、IEA分析")

# P17: 全球补贴政策地图
add_two_section_slide(prs,
    "补贴从'压低化石能源价格'转向'投资可再生能源'",
    "历史性政策拐点已出现——各国补贴模式正在不可逆转变",
    "【旧模式：化石能源补贴】\n• 压低汽油/柴油价格\n• 天然气价格管制\n• 电价上限\n• 财政负担重\n• 不可持续\n\n代表国家:\n• 巴基斯坦（财政困难）\n• 孟加拉国（难以维持）",
    "【新模式：可再生能源投资】\n• 日本: $31亿夏季电力补贴\n• 德国: 工业电价补贴\n• 挪威: 居民电价上限\n• 中国: 千亿级新能源投资/年\n• 韩国: 电价稳定基金\n\n关键趋势:\n• 泰国已放弃汽油价格上限\n• 转向市场化+新能源\n• 补贴与投资并重",
    "数据来源：IEA, 各国财政部/能源部公告")

# P18: CTA
add_cta_slide(prs,
    "窗口期不会永远打开",
    [
        "建立LNG依存国新能源机会地图",
        "组建海外业务专项小组",
        "启动泰国、日本初步接触"
    ],
    "Q2内完成评估 → Q3泰国试点 → Q4日本首单 → 2027规模化",
    "机密 | 仅供内部战略讨论使用")

# Save
prs.save(output_path)
print(f"✅ Deck saved to: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
